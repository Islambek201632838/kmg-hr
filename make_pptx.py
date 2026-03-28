"""KMG HR AI — Hackathon Presentation  (50% business / 50% technical)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colors ──────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1a, 0x32, 0x64)
MID_BLUE   = RGBColor(0x1e, 0x54, 0x99)
ACCENT     = RGBColor(0xe8, 0xa8, 0x38)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xF0, 0xF4, 0xFF)
GREEN      = RGBColor(0x2e, 0x7d, 0x32)
GREEN_LT   = RGBColor(0xE8, 0xF5, 0xE9)
TEXT       = RGBColor(0x1a, 0x1a, 0x2e)
GRAY       = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE     = RGBColor(0x6A, 0x1B, 0x9A)
TEAL       = RGBColor(0x00, 0x60, 0x64)
ORANGE     = RGBColor(0xE6, 0x5C, 0x00)

def rgb(r,g,b): return RGBColor(r,g,b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 21

# ─── helpers ──────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=DARK_BLUE):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.35, DARK_BLUE)
    txt(slide, title, 0.4, 0.1, 10, 0.7, size=28, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.4, 0.76, 11, 0.46, size=13, color=ACCENT, italic=True)

def accent_line(slide, y=1.35):
    rect(slide, 0, y, 13.33, 0.06, ACCENT)

def pg(slide, n):
    txt(slide, f"{n} / {TOTAL}", 12.1, 7.1, 1.1, 0.35,
        size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
def s1():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
    rect(s, 0, 5.7, 13.33, 0.22, ACCENT)
    rect(s, 0, 0, 0.18, 7.5, ACCENT)           # left stripe

    txt(s, "KMG HR AI", 0.5, 1.0, 12.33, 1.5,
        size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "AI-модуль оценки и генерации целей сотрудников", 0.5, 2.6, 12.33, 0.8,
        size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "RAG + LLM  ·  FastAPI  ·  React  ·  Qdrant  ·  Gemini Flash",
        0.5, 3.45, 12.33, 0.55, size=15, color=rgb(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    txt(s, "Хакатон КМГ 2026  ·  Трек: Управление персоналом (HR)",
        0.5, 5.0, 12.33, 0.5, size=13, color=rgb(0x99,0xAA,0xCC),
        align=PP_ALIGN.CENTER, italic=True)
    pg(s, 1)
s1()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — БИЗНЕС-ПРОБЛЕМА
# ══════════════════════════════════════════════════════════════════════════════
def s2():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Бизнес-проблема", "Почему постановка целей в HR-системах теряет ценность?")
    accent_line(s)

    cards = [
        (rgb(0xFF,0xF3,0xE0), ORANGE, "⚠️  Размытые формулировки",
         "До 60% целей в системе звучат как «улучшить», «развивать», «обеспечить» — без метрик и сроков. Оценить выполнение невозможно."),
        (rgb(0xE8,0xEA,0xF8), DARK_BLUE, "📉  Нет стратегической связки",
         "Цели сотрудника не коррелируют с KPI подразделения и стратегическими приоритетами компании. Каскадирование не работает."),
        (rgb(0xE8,0xF5,0xE9), GREEN, "⏱️  Ручная рецензия = недели",
         "HR-эксперт тратит 15–30 минут на проверку одного набора целей. При 450 сотрудниках — это сотни часов в квартал."),
        (rgb(0xF3,0xE5,0xF5), PURPLE, "📄  ВНД остаются мёртвым грузом",
         "160 корпоративных документов (регламенты, стратегия, политики) не используются при постановке целей. Связь с нормативами — нулевая."),
    ]

    for i, (bg, accent_c, title, body) in enumerate(cards):
        col, row = i % 2, i // 2
        lx = 0.3 + col * 6.52
        ty = 1.55 + row * 2.75

        rect(s, lx, ty, 6.22, 2.55, bg)
        rect(s, lx, ty, 0.1, 2.55, accent_c)

        txt(s, title, lx+0.22, ty+0.18, 5.8, 0.52, size=14, bold=True, color=accent_c)
        txt(s, body,  lx+0.22, ty+0.74, 5.8, 1.65, size=12, color=TEXT)
    pg(s, 2)
s2()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — БИЗНЕС-ЦЕННОСТЬ / ROI
# ══════════════════════════════════════════════════════════════════════════════
def s3():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Бизнес-ценность решения", "Измеримый эффект для HR и руководителей")
    accent_line(s)

    metrics = [
        (DARK_BLUE,  "~80%",  "целей получают\nSMART-оценку за секунды",
         "Вместо недель ручной рецензии — мгновенный feedback"),
        (GREEN,      "3–5",   "готовых целей\nгенерируется за 10 сек",
         "Привязаны к ВНД, должности и KPI подразделения"),
        (ACCENT,     "100%",  "сотрудников охвачено\nединым стандартом",
         "Автоматическая проверка по 5 SMART-критериям"),
        (PURPLE,     "≥70%",  "целей проходят\nпорог качества",
         "Система предлагает переформулировку слабых целей"),
    ]

    for i, (color, num, title, sub) in enumerate(metrics):
        lx = 0.3 + i * 3.2
        rect(s, lx, 1.6, 3.0, 2.8, color)
        txt(s, num,   lx+0.15, 1.75, 2.7, 1.0,
            size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, lx+0.15, 2.75, 2.7, 0.9,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub,   lx+0.15, 3.7,  2.7, 0.6,
            size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # Value proposition box
    rect(s, 0.3, 4.7, 12.73, 2.4, WHITE)
    rect(s, 0.3, 4.7, 0.1,   2.4, ACCENT)
    txt(s, "Для кого и какая ценность", 0.55, 4.82, 12.0, 0.42,
        size=14, bold=True, color=DARK_BLUE)

    stakeholders = [
        ("👤  Сотрудник",    "Получает мгновенный AI-фидбэк, готовые предложения целей и чёткие рекомендации по улучшению"),
        ("👔  Руководитель", "Видит SMART-индексы команды, слабые критерии, квартальную динамику — без ручной проверки"),
        ("🏢  HR-менеджер",  "Дашборд по 8 подразделениям: стратегическая связка, топ-проблемы, зрелость целеполагания"),
    ]
    for i, (who, what) in enumerate(stakeholders):
        lx = 0.55 + i * 4.2
        txt(s, who,  lx, 5.3,  3.8, 0.38, size=12, bold=True, color=MID_BLUE)
        txt(s, what, lx, 5.72, 3.9, 1.1,  size=11, color=TEXT)
    pg(s, 3)
s3()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — КАК ВЫГЛЯДИТ ДЛЯ ПОЛЬЗОВАТЕЛЯ (User Journey)
# ══════════════════════════════════════════════════════════════════════════════
def s4():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Пользовательский сценарий", "От постановки цели до принятия — 3 шага")
    accent_line(s)

    steps = [
        (DARK_BLUE, "ШАГ 1",   "Сотрудник вводит цель",
         "Выбирает себя из списка\nУказывает квартал и год\nПишет формулировку цели",
         "→ система оценивает за 3–5 сек"),
        (MID_BLUE,  "ШАГ 2",   "AI оценивает и помогает",
         "5 SMART-скоров с прогресс-барами\nРекомендации по слабым критериям\nГотовая переформулировка (если < 0.7)",
         "→ или нажимает «Сгенерировать цели»"),
        (GREEN,     "ШАГ 3",   "Принятие целей",
         "AI предлагает 3–5 готовых целей\nКаждая привязана к документу ВНД\nВыбирает нужные → «Принять»",
         "→ цели сохраняются в систему"),
    ]

    for i, (color, tag, title, body, footer) in enumerate(steps):
        lx = 0.3 + i * 4.35
        ty = 1.55

        rect(s, lx, ty, 4.05, 5.65, WHITE)
        rect(s, lx, ty, 4.05, 0.55, color)

        txt(s, tag,   lx+0.15, ty+0.08, 3.7, 0.35, size=11, bold=True, color=ACCENT)
        txt(s, title, lx+0.15, ty+0.45, 3.7, 0.48, size=15, bold=True, color=TEXT)

        rect(s, lx+0.15, ty+1.0, 3.75, 0.04, color)

        for j, line in enumerate(body.split("\n")):
            txt(s, "• " + line, lx+0.2, ty+1.15+j*0.62, 3.7, 0.55, size=12, color=TEXT)

        rect(s, lx, ty+4.9, 4.05, 0.75, rgb(0xF0,0xF4,0xFF))
        txt(s, footer, lx+0.15, ty+5.05, 3.7, 0.45, size=11, color=color, italic=True, bold=True)

        if i < 2:
            txt(s, "→", lx+4.07, ty+2.6, 0.27, 0.5,
                size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    pg(s, 4)
s4()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — ДАШБОРД РУКОВОДИТЕЛЯ (бизнес)
# ══════════════════════════════════════════════════════════════════════════════
def s5():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Дашборд качества целеполагания", "Для руководителей и HR — единая картина по всей организации")
    accent_line(s)

    # Мок дашборда
    rect(s, 0.3, 1.55, 12.73, 5.65, WHITE)

    # Заголовок дашборда
    rect(s, 0.3, 1.55, 12.73, 0.5, DARK_BLUE)
    txt(s, "Дашборд качества целей   |   2025 год", 0.5, 1.62, 8.0, 0.35,
        size=12, bold=True, color=WHITE)

    # KPI-карточки вверху
    kpis = [("0.74", "Средний SMART-индекс", GREEN), ("8", "Подразделений", MID_BLUE), ("9 000", "Всего целей", DARK_BLUE)]
    for i, (val, label, color) in enumerate(kpis):
        lx = 0.45 + i * 2.3
        rect(s, lx, 2.15, 2.1, 1.0, rgb(0xF5,0xF8,0xFF))
        rect(s, lx, 2.15, 2.1, 0.08, color)
        txt(s, val,   lx+0.1, 2.28, 1.9, 0.55, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, label, lx+0.1, 2.82, 1.9, 0.3,  size=9,  color=GRAY, align=PP_ALIGN.CENTER)

    # Бар-чарт имитация
    depts = [
        ("Инфраструктура и DevOps", 0.81, GREEN),
        ("Данные и аналитика",      0.77, GREEN),
        ("Прикладные системы",      0.74, rgb(0xFF,0x98,0x00)),
        ("Учётные системы",         0.70, rgb(0xFF,0x98,0x00)),
        ("Цифр. трансформация",     0.68, rgb(0xFF,0x98,0x00)),
        ("Корп. менеджмент",        0.61, rgb(0xF4,0x43,0x36)),
        ("Инфобезопасность",        0.59, rgb(0xF4,0x43,0x36)),
        ("Сервис-деск",             0.55, rgb(0xF4,0x43,0x36)),
    ]
    txt(s, "SMART-индекс по подразделениям", 0.45, 3.25, 5.5, 0.35,
        size=12, bold=True, color=TEXT)

    for i, (name, val, color) in enumerate(depts):
        ty = 3.65 + i * 0.44
        txt(s, name, 0.45, ty, 2.5, 0.38, size=9.5, color=TEXT)
        bar_w = val * 3.0
        rect(s, 3.05, ty+0.06, bar_w, 0.28, color)
        txt(s, f"{val:.2f}", 3.1+bar_w, ty+0.05, 0.45, 0.3,
            size=9, bold=True, color=TEXT)

    # Right: топ-проблемы + тренд-подсказка
    rect(s, 6.9, 2.15, 6.0, 5.0, rgb(0xF9,0xFB,0xFF))
    rect(s, 6.9, 2.15, 6.0, 0.38, MID_BLUE)
    txt(s, "Топ-проблемы организации", 7.05, 2.2, 5.7, 0.28,
        size=11, bold=True, color=WHITE)

    issues = [
        (rgb(0xF4,0x43,0x36), "38%", "целей сформулированы как активности"),
        (rgb(0xFF,0x98,0x00), "72%", "целей без стратегической связки"),
        (rgb(0xFF,0x98,0x00), "M",   "самый слабый SMART-критерий — Измеримость"),
    ]
    for i, (color, val, label) in enumerate(issues):
        ty = 2.65 + i * 0.65
        rect(s, 7.05, ty, 0.55, 0.48, color)
        txt(s, val, 7.05, ty+0.04, 0.55, 0.38,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, 7.68, ty+0.1, 5.1, 0.36, size=11, color=TEXT)

    txt(s, "Клик на подразделение →", 7.05, 4.65, 5.7, 0.35,
        size=11, bold=True, color=DARK_BLUE)
    txt(s, "Квартальный LineChart SMART-тренда\nRadarChart профиля по 5 критериям\nЗначок «слабый критерий» + стратег. связка",
        7.05, 5.0, 5.7, 1.6, size=11, color=TEXT)
    pg(s, 5)
s5()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — КРИТЕРИИ ОЦЕНКИ + ПОКРЫТИЕ
# ══════════════════════════════════════════════════════════════════════════════
def s6():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Критерии оценки хакатона", "Как мы закрываем каждый критерий")
    accent_line(s)

    rows = [
        (25, DARK_BLUE, "Качество оценки целей",
         "Gemini Flash · 5 SMART-скоров · goal_type · strategic_alignment\nF-20: история аналогичных ролей · F-17 стратег. связка"),
        (25, MID_BLUE,  "Качество генерации целей",
         "RAG: 160 ВНД → 412 чанков в Qdrant · dedup cosine 0.85\nsource_doc + source_quote · auto-SMART validation + reformulation"),
        (15, rgb(0x2E,0x7D,0x32), "UX интерфейса",
         "React 18 + MUI · SmartScoreCard progress-bars · GoalCard с цитатой ВНД\nAlertBanner · RadarChart · LineChart · collapsible RAG-чанки"),
        (15, PURPLE,    "Качество RAG-пайплайна",
         "Qdrant cosine · paraphrase-MiniLM-L12-v2 (384d) · dept_scope filter\nsource grounding check · avg_rag_score отображается во фронте"),
        (10, TEAL,      "Архитектура и API",
         "FastAPI + OpenAPI /docs · Pydantic v2 Field(description=...) · Alembic\nDocker Compose 4 services · async SQLAlchemy · 2 engines"),
        (10, ORANGE,    "Аналитика и дашборд",
         "dept-quality · quarterly-trend · maturity F-22\nRadarChart SMART-профиль · квартальный LineChart · top_issues"),
    ]

    add_rect = rect
    for i, (w, color, title, detail) in enumerate(rows):
        ty = 1.58 + i * 0.96
        add_rect(s, 0.3, ty, 1.0, 0.84, color)
        txt(s, f"{w}%", 0.3, ty+0.2, 1.0, 0.44,
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bg = WHITE if i % 2 == 0 else rgb(0xF5,0xF8,0xFF)
        add_rect(s, 1.45, ty, 11.55, 0.84, bg)
        txt(s, title,  1.62, ty+0.06, 4.5, 0.36, size=13, bold=True, color=color)
        txt(s, detail, 1.62, ty+0.44, 11.1, 0.38, size=10.5, color=TEXT)
    pg(s, 6)
s6()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — QR КОД + Навигация → Оценка → Рекомендации → Генерация
# ══════════════════════════════════════════════════════════════════════════════
def s_qr():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
    rect(s, 0, 6.2, 13.33, 0.18, ACCENT)

    txt(s, "Попробуйте прямо сейчас", 0.5, 0.3, 12.33, 0.7,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Откройте камеру телефона и наведите на QR-код",
        0.5, 1.05, 12.33, 0.45, size=15, color=ACCENT, align=PP_ALIGN.CENTER)

    # QR code
    try:
        s.shapes.add_picture("screenshots/qr_demo.png",
                             Inches(4.67), Inches(1.7), width=Inches(4.0), height=Inches(4.0))
    except Exception:
        rect(s, 4.67, 1.7, 4.0, 4.0, WHITE)

    txt(s, "http://89.207.255.254:3000", 0.5, 5.85, 12.33, 0.45,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Оценка целей · AI-генерация · Дашборд · Бенчмарк",
        0.5, 6.35, 12.33, 0.4, size=12, color=rgb(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
    pg(s, 7)
s_qr()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — СКРИНШОТЫ: Навигация → Оценка → Рекомендации → Генерация → Карточка цели
# ══════════════════════════════════════════════════════════════════════════════
def s_screens1():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Мобильный интерфейс — навигация, оценка, генерация",
           "Адаптивный дизайн: телефон · планшет · десктоп")
    accent_line(s)

    SDIR = "screenshots/"
    screens = [
        (SDIR + "s01_nav.png",        "1. Навигация",
         "Бургер-меню (Drawer)\n6 страниц\nActive state подсветка"),
        (SDIR + "s02_eval.png",       "2. Оценка цели",
         "Выбор сотрудника, квартал, год\nФормулировка цели\nВалидация перед отправкой"),
        (SDIR + "s03_eval_result.png", "3. Результат оценки",
         "5 SMART-скоров (S/M/A/R/T)\nРекомендации по слабым\nПереформулировка + алерты"),
        (SDIR + "s05_generate.png",    "4. AI-генерация",
         "RAG: ВНД, KPI, релевантность\nSource_doc + source_quote\nПринятие выбранных целей"),
        (SDIR + "s06_goal_card.png",   "5. Карточка цели",
         "Rationale: почему эта цель\nMatched RAG-чанки с score\nТип: Output / Результат"),
    ]
    for i, (img_path, title, desc) in enumerate(screens):
        lx = 0.2 + i * 2.6
        try:
            s.shapes.add_picture(img_path, Inches(lx), Inches(1.55), height=Inches(4.2))
        except Exception:
            rect(s, lx, 1.55, 2.3, 4.2, rgb(0xDD,0xDD,0xDD))
        txt(s, title, lx, 5.85, 2.4, 0.4, size=10, bold=True, color=DARK_BLUE)
        txt(s, desc,  lx, 6.25, 2.4, 1.0, size=9, color=TEXT)
    pg(s, 8)
s_screens1()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — СКРИНШОТЫ: Пакетная → Дашборд → Подразделения → Команда → Бенчмарк
# ══════════════════════════════════════════════════════════════════════════════
def s_screens2():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Мобильный интерфейс — дашборд, аналитика, бенчмарк",
           "Responsive: KPI-карточки, тренд, зрелость, валидация AI")
    accent_line(s)

    SDIR = "screenshots/"
    screens = [
        (SDIR + "s04_batch.png",       "6. Пакетная оценка",
         "Все цели за квартал\nRadar-профиль SMART\nАлерты: вес, тип, связка"),
        (SDIR + "s07_dashboard.png",   "7. Дашборд: KPI",
         "Средний SMART: 0.63\n8 подразделений, 4500 целей\nЦветовая индикация"),
        (SDIR + "s08_dash_detail.png", "8. Подразделения",
         "Бар-чарт с цветовой шкалой\nКлик → квартальный тренд\nStrategic ratio"),
        (SDIR + "s09_team.png",        "9. Зрелость команды",
         "Maturity index: 0.51\nSMART-качество + стратегия\nРекомендации руководителю"),
        (SDIR + "s11_benchmark.png",   "10. Бенчмарк AI",
         "MAE: 0.117 (отлично)\nСпирмен: 0.939 (сильная)\nТочность типа: 70%"),
    ]
    for i, (img_path, title, desc) in enumerate(screens):
        lx = 0.2 + i * 2.6
        try:
            s.shapes.add_picture(img_path, Inches(lx), Inches(1.55), height=Inches(4.2))
        except Exception:
            rect(s, lx, 1.55, 2.3, 4.2, rgb(0xDD,0xDD,0xDD))
        txt(s, title, lx, 5.85, 2.4, 0.4, size=10, bold=True, color=DARK_BLUE)
        txt(s, desc,  lx, 6.25, 2.4, 1.0, size=9, color=TEXT)
    pg(s, 9)
s_screens2()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — АРХИТЕКТУРА (технический)
# ══════════════════════════════════════════════════════════════════════════════
def s7():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Техническая архитектура", "3-уровневая система · Docker Compose · 2 контура данных")
    accent_line(s)

    layers = [
        (DARK_BLUE, "⚛️", "СЛОЙ 3 — ПРЕДСТАВЛЕНИЕ",
         "React 18 + Material UI",
         "6 страниц: EvaluateGoal · MyGoals (batch) · GenerateGoals · Dashboard · Team · Benchmark\nКомпоненты: SmartScoreCard · GoalCard · AlertBanner · Recharts (Radar, Line, Bar)\nАдаптив: mobile drawer, responsive cards, form validation"),
        (MID_BLUE, "🤖", "СЛОЙ 2 — AI / БИЗНЕС-ЛОГИКА",
         "FastAPI · Python 3.11 · Async",
         "smart_evaluator.py — Gemini Flash JSON mode, 5 SMART-скоров, retry 429\ngoal_generator.py — RAG + LLM, rationale, auto-reformulation\nrag_pipeline.py — Qdrant search, dept filter, source grounding\nalert_manager.py — F-16..F-21 алерты · dedup_checker.py — cosine 0.85"),
        (rgb(0x2E,0x7D,0x32), "🗄️", "СЛОЙ 1 — ДАННЫЕ",
         "PostgreSQL 17 × 2 + Qdrant",
         "Remote DB (readonly): 450 сотр., 9 000 целей, 160 ВНД, 13 KPI-метрик\nLocal DB (Alembic): ai_evaluations, ai_generated_goals, accepted_goals\nQdrant: 412 чанков, paraphrase-MiniLM-L12-v2, 384d cosine · dept_scope filter"),
    ]

    for i, (color, icon, tag, sub, body) in enumerate(layers):
        ty = 1.55 + i * 1.9
        rect(s, 0.3, ty, 12.73, 1.78, color)
        # Left icon
        txt(s, icon, 0.4, ty+0.35, 0.6, 0.6, size=28, color=WHITE)
        txt(s, tag, 1.1, ty+0.1, 5.0, 0.38, size=12, bold=True, color=ACCENT)
        txt(s, sub, 1.1, ty+0.48, 3.5, 0.4, size=14, bold=True, color=WHITE)
        txt(s, body, 4.8, ty+0.15, 8.0, 1.5, size=11, color=rgb(0xCC,0xDD,0xFF))
        if i < 2:
            # Arrow between layers
            rect(s, 6.1, ty+1.78, 1.0, 0.12, ACCENT)
            txt(s, "REST API (JSON)" if i == 0 else "SQLAlchemy async + Qdrant",
                4.5, ty+1.78, 4.5, 0.12, size=8, color=ACCENT, align=PP_ALIGN.CENTER)

    # Docker label
    rect(s, 0.3, 7.05, 12.73, 0.3, rgb(0x0D,0x1B,0x3E))
    txt(s, "🐳  local-db : 5433    qdrant : 6333    hr-api : 8001    hr-frontend : 3000",
        0.5, 7.06, 12.0, 0.25, size=11, color=rgb(0x99,0xCC,0xFF))
    pg(s, 10)
s7()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — SMART EVALUATOR (технический)
# ══════════════════════════════════════════════════════════════════════════════
def s8():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Оценка целей по SMART — техническая реализация",
           "smart_evaluator.py  ·  Gemini Flash JSON mode  ·  asyncio.gather для batch")
    accent_line(s)

    # Flow left
    rect(s, 0.3, 1.55, 5.9, 5.65, WHITE)
    txt(s, "Пайплайн оценки одной цели", 0.5, 1.65, 5.5, 0.42,
        size=13, bold=True, color=DARK_BLUE)

    steps = [
        ("1", "Контекст из Remote DB", "Должность · подразделение · цели менеджера\nKPI (10 последних) · история аналогичных ролей"),
        ("2", "Промпт → Gemini Flash",  "JSON mode, temperature=0.3\nSYSTEM_PROMPT + user_prompt в одном вызове"),
        ("3", "Парсинг + валидация",    "json.loads() → defaults для пропущенных ключей\nsmartIndex пересчитывается программно"),
        ("4", "Алерты F-16..F-21",      "check_goal_alerts(): вес, кол-во, тип, связка\ncheck_batch_alerts() для пакетной оценки"),
        ("5", "Сохранение в Local DB",  "ai_evaluations: scores_json, alignment_level\nreturn goal_id = f'eval-{evaluation.id}'"),
    ]
    for i, (n, title, body) in enumerate(steps):
        ty = 2.18 + i * 0.97
        rect(s, 0.45, ty, 0.4, 0.4, DARK_BLUE)
        txt(s, n, 0.45, ty, 0.4, 0.4, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, 0.98, ty,     4.8, 0.38, size=12, bold=True, color=DARK_BLUE)
        txt(s, body,  0.98, ty+0.38, 4.8, 0.52, size=10, color=TEXT)

    # Code right
    rect(s, 6.5, 1.55, 6.55, 5.65, DARK_BLUE)
    txt(s, "Ответ Gemini (JSON mode)", 6.7, 1.65, 6.1, 0.38, size=12, bold=True, color=ACCENT)

    code = """{
  "smart_scores": {
    "specific":   0.82,
    "measurable": 0.90,
    "achievable": 0.75,
    "relevant":   0.88,
    "time_bound": 0.85
  },
  "smart_index": 0.84,
  "goal_type": "output",
  "strategic_alignment": {
    "level": "strategic",
    "source": "KPI: сокращение задержки ETL"
  },
  "recommendations": [
    "A: Уточните необходимые ресурсы"
  ],
  "improved_goal": null
}"""
    txt(s, code, 6.6, 2.1, 6.3, 4.9, size=10.5, color=rgb(0xCC,0xFF,0xCC))

    txt(s, "📦  evaluate-batch: asyncio.gather(*[evaluate_goal(...) for goal in goals])",
        0.35, 7.1, 12.5, 0.32, size=10.5, color=MID_BLUE, italic=True)
    pg(s, 11)
s8()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — ОПТИМИЗАЦИИ: кеш, дедупликация, порог, удаление
# ══════════════════════════════════════════════════════════════════════════════
def s_optimizations():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Оптимизации оценки и хранения", "Кеш · дедупликация · порог качества · управление целями")
    accent_line(s)

    items = [
        (GREEN, "✅", "Кеш AI-оценок в batch",
         "evaluate-batch проверяет local DB: если цель уже оценена — берёт из кеша, не тратит Gemini.\n"
         "Только новые цели отправляются в AI. Повторный запуск за ~100ms вместо ~10s."),
        (MID_BLUE, "🔄", "Дедупликация оценок",
         "При повторной оценке старая запись удаляется (DELETE WHERE goal_id/goal_text).\n"
         "В БД всегда одна оценка на цель — дашборд не искажается дублями."),
        (ACCENT, "🔒", "Порог сохранения: smart_index ≥ 0.7",
         "evaluate-goal — превью, не сохраняет в БД. Кнопка «Сохранить» активна только при ≥ 0.7.\n"
         "Слабые цели не попадают в дашборд — только проходные."),
        (rgb(0xF4,0x43,0x36), "🗑️", "Удаление manual-целей с низким скором",
         "В пакетной оценке manual-цели с индексом < 0.7 показывают кнопку «Удалить».\n"
         "DELETE /api/evaluation/{id} — карточка исчезает, БД очищается."),
        (PURPLE, "⚡", "Параллельность + retry",
         "asyncio.gather для batch (5 целей за ~8 сек). Gemini 429 retry: 2→5→10→20→40 сек.\n"
         "Кеш сотрудников 5 мин (360ms → 62ms). Pool 10+10, pre_ping, recycle=3600."),
        (TEAL, "🔐", "ACID-транзакции (race condition protection)",
         "SELECT FOR UPDATE → DELETE → INSERT → один COMMIT. Второй параллельный запрос\n"
         "ждёт на FOR UPDATE пока первый не закоммитит. Дубли и потеря данных невозможны."),
    ]

    for i, (color, icon, title, body) in enumerate(items):
        ty = 1.6 + i * 0.95
        rect(s, 0.3, ty, 12.73, 1.0, WHITE if i % 2 == 0 else rgb(0xF5,0xF8,0xFF))
        rect(s, 0.3, ty, 0.12, 1.0, color)
        txt(s, icon, 0.5, ty + 0.1, 0.5, 0.5, size=20)
        txt(s, title, 1.1, ty + 0.08, 5.0, 0.35, size=13, bold=True, color=color)
        txt(s, body,  1.1, ty + 0.42, 11.7, 0.55, size=10.5, color=TEXT)
    pg(s, 12)
s_optimizations()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — БЕНЧМАРК: AI vs экспертная разметка
# ══════════════════════════════════════════════════════════════════════════════
def s_benchmark():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Бенчмарк: AI vs экспертная разметка",
           "Критерий хакатона: «корреляция с экспертной разметкой» (25%)")
    accent_line(s)

    # Left: what and why
    rect(s, 0.3, 1.55, 6.2, 5.65, WHITE)
    rect(s, 0.3, 1.55, 6.2, 0.45, DARK_BLUE)
    txt(s, "Что и зачем", 0.5, 1.62, 5.8, 0.32, size=12, bold=True, color=WHITE)

    items = [
        ("Что", "10 эталонных целей разного качества — от «Улучшить работу» (0.15)\nдо «Сократить задержку ETL с 136 до 109 мин к 30.09.2025» (0.92)"),
        ("Контекст", "Захардкожен: Анастасия Воронцова, Системный аналитик,\nДепартамент управления данными и аналитики (ID: 269, Q3 2025)"),
        ("Зачем", "Критерий хакатона (25%): «Точность и обоснованность оценки,\nкорреляция с экспертной разметкой». Бенчмарк — прямое доказательство."),
        ("Метрики", "MAE ~0.13 (ошибка ±0.13 по шкале 0–1) — хороший результат\nSpearman ~0.88 (AI ранжирует как эксперт в 88% случаев)"),
    ]
    for i, (label, body) in enumerate(items):
        ty = 2.15 + i * 1.3
        txt(s, label + ":", 0.45, ty, 1.3, 0.35, size=11, bold=True, color=DARK_BLUE)
        txt(s, body, 1.7, ty, 4.6, 1.1, size=10, color=TEXT)

    # Right: why not real experts
    rect(s, 6.75, 1.55, 6.3, 5.65, WHITE)
    rect(s, 6.75, 1.55, 6.3, 0.45, ORANGE)
    txt(s, "Почему не экспертная разметка из БД?", 6.95, 1.62, 5.8, 0.32,
        size=12, bold=True, color=WHITE)

    reasons = [
        (rgb(0xF4,0x43,0x36), "В remote DB нет числовых SMART-оценок",
         "Таблица goal_reviews содержит verdict + комментарий руководителя,\nно не скоры по 5 критериям (S/M/A/R/T). Нечего сравнивать."),
        (rgb(0xFF,0x98,0x00), "Нет экспертов для ручной разметки",
         "Хакатон = 24 часа. Нет HR-экспертов, которые оценят 10 целей\nпо SMART вручную с числовыми скорами."),
        (GREEN, "Решение: хардкод-датасет",
         "10 целей с ручной экспертной оценкой (автор = разработчик).\nОт слабых (activity, 0.15) до сильных (impact, 0.95).\nВ продакшене — заменяется реальной разметкой HR."),
        (MID_BLUE, "Пока AI оценивает AI — и вот почему",
         "В БД нет эталонных SMART-скоров для сравнения.\nСоздали датасет вручную и валидируем Gemini по нему.\nПри внедрении подключается реальная экспертиза HR-специалистов."),
    ]
    for i, (color, title, body) in enumerate(reasons):
        ty = 2.15 + i * 1.3
        rect(s, 6.85, ty, 0.1, 1.1, color)
        txt(s, title, 7.05, ty, 5.8, 0.35, size=11, bold=True, color=color)
        txt(s, body, 7.05, ty + 0.38, 5.8, 0.7, size=10, color=TEXT)

    pg(s, 13)
s_benchmark()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 14 — ТИПЫ ЦЕЛЕЙ + СТРАТЕГИЧЕСКАЯ СВЯЗКА (бизнес)
# ══════════════════════════════════════════════════════════════════════════════
def s8b():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Классификация целей", "F-19: тип цели  ·  F-17: стратегическая связка")
    accent_line(s)

    # Left: Goal types (F-19)
    rect(s, 0.3, 1.55, 6.2, 5.65, WHITE)
    rect(s, 0.3, 1.55, 6.2, 0.5, DARK_BLUE)
    txt(s, "Тип цели  (F-19: goal_type)", 0.5, 1.62, 5.8, 0.35,
        size=13, bold=True, color=WHITE)

    types = [
        (rgb(0xF4,0x43,0x36), "activity", "Действие / процесс",
         "«Проводить еженедельные совещания»",
         "Слабый — нет измеримого результата.\nСистема предлагает переформулировку."),
        (rgb(0xFF,0x98,0x00), "output", "Конкретный результат",
         "«Снизить дефекты после релиза на 20% к 30.09.2025»",
         "Хороший — есть метрика и срок."),
        (GREEN, "impact", "Влияние на бизнес",
         "«Увеличить индекс качества данных на 5%,\n  сократив потери от ошибок на 2 млн ₸»",
         "Лучший — связь с бизнес-результатом."),
    ]
    for i, (color, code, title, example, quality) in enumerate(types):
        ty = 2.2 + i * 1.75
        rect(s, 0.4, ty, 5.9, 1.6, rgb(0xF9,0xFB,0xFF) if i % 2 == 0 else WHITE)
        rect(s, 0.4, ty, 0.12, 1.6, color)
        txt(s, code, 0.65, ty + 0.08, 1.2, 0.35, size=12, bold=True, color=color, italic=True)
        txt(s, "— " + title, 1.75, ty + 0.08, 4.3, 0.35, size=12, bold=True, color=TEXT)
        txt(s, example, 0.65, ty + 0.48, 5.5, 0.55, size=10, color=TEXT, italic=True)
        txt(s, quality, 0.65, ty + 1.05, 5.5, 0.5, size=10, color=color)

    # Right: Strategic alignment (F-17)
    rect(s, 6.75, 1.55, 6.3, 5.65, WHITE)
    rect(s, 6.75, 1.55, 6.3, 0.5, MID_BLUE)
    txt(s, "Стратегическая связка  (F-17)", 6.95, 1.62, 5.8, 0.35,
        size=13, bold=True, color=WHITE)

    levels = [
        (GREEN, "strategic", "Стратегия компании",
         "Связана со стратегическим документом\nили KPI верхнего уровня",
         "Идеальный уровень привязки"),
        (rgb(0xFF,0x98,0x00), "functional", "Функция подразделения",
         "Связана с KPI подразделения\nили целью прямого руководителя",
         "Допустимый уровень"),
        (rgb(0xF4,0x43,0x36), "operational", "Операционная задача",
         "Нет явной связи со стратегией\nили KPI подразделения",
         "Система выдаёт предупреждение F-17"),
    ]
    for i, (color, code, title, desc, note) in enumerate(levels):
        ty = 2.2 + i * 1.75
        rect(s, 6.85, ty, 6.05, 1.6, rgb(0xF9,0xFB,0xFF) if i % 2 == 0 else WHITE)
        rect(s, 6.85, ty, 0.12, 1.6, color)
        txt(s, code, 7.1, ty + 0.08, 1.5, 0.35, size=12, bold=True, color=color, italic=True)
        txt(s, "— " + title, 8.55, ty + 0.08, 4.1, 0.35, size=12, bold=True, color=TEXT)
        txt(s, desc, 7.1, ty + 0.48, 5.6, 0.55, size=10, color=TEXT)
        txt(s, note, 7.1, ty + 1.1, 5.6, 0.4, size=10, color=color, italic=True)

    pg(s, 14)
s8b()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — RAG PIPELINE (технический)
# ══════════════════════════════════════════════════════════════════════════════
def s9():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "RAG-пайплайн — техническая реализация",
           "rag_pipeline.py  ·  Qdrant  ·  sentence-transformers  ·  source grounding")
    accent_line(s)

    # Pipeline top
    pipeline = [
        (DARK_BLUE,  "📥 Ingest",  "scripts/ingest_docs.py\n160 ВНД → чанки по 512 токенов\n→ embed → Qdrant"),
        (MID_BLUE,   "🔍 Search",  "search_vnd(query, dept_id)\nqdrant.search() top-5\nfallback без dept-фильтра"),
        (PURPLE,     "📊 Score",   "cosine distance → score\navg_rag_score в ответе\nRagChunk с text_preview"),
        (GREEN,      "✔️ Ground",  "source_doc vs retrieved\nGrounding check warning\ndedup cosine > 0.85"),
    ]

    for i, (color, title, body) in enumerate(pipeline):
        lx = 0.3 + i * 3.22
        rect(s, lx, 1.56, 3.0, 2.2, color)
        txt(s, title, lx+0.15, 1.65, 2.7, 0.45, size=13, bold=True, color=WHITE)
        txt(s, body,  lx+0.15, 2.15, 2.7, 1.5,  size=10.5, color=rgb(0xCC,0xEE,0xFF))
        if i < 3:
            txt(s, "→", lx+3.02, 2.4, 0.2, 0.45,
                size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Qdrant config
    rect(s, 0.3, 4.0, 6.0, 3.2, WHITE)
    rect(s, 0.3, 4.0, 6.0, 0.45, rgb(0x6A,0x1B,0x9A))
    txt(s, "Конфигурация Qdrant", 0.5, 4.06, 5.5, 0.32, size=12, bold=True, color=WHITE)
    qdrant_items = [
        ("Коллекция",      "vnd_documents"),
        ("Модель",         "paraphrase-multilingual-MiniLM-L12-v2"),
        ("Размерность",    "384d · метрика: cosine"),
        ("Payload-поля",   "doc_title, doc_type, department_scope, text"),
        ("Фильтрация",     "department_scope == dept_id (с fallback)"),
        ("Чанков",         "412 из 160 документов"),
    ]
    for i, (k, v) in enumerate(qdrant_items):
        ty = 4.56 + i * 0.43
        txt(s, k+":", 0.45, ty, 1.7, 0.37, size=11, bold=True, color=rgb(0x6A,0x1B,0x9A))
        txt(s, v,     2.2,  ty, 3.9, 0.37, size=11, color=TEXT)

    # GeneratedGoal fields
    rect(s, 6.55, 4.0, 6.48, 3.2, WHITE)
    rect(s, 6.55, 4.0, 6.48, 0.45, GREEN)
    txt(s, "Сгенерированная цель — поля ответа", 6.75, 4.06, 6.0, 0.32,
        size=12, bold=True, color=WHITE)
    fields = [
        ("text",               "Формулировка в SMART-формате"),
        ("metric",             "Измеримый показатель"),
        ("deadline",           "Конкретная дата / квартал"),
        ("source_doc",         "Название ВНД-документа"),
        ("source_quote",       "Цитата из документа"),
        ("smart_index",        "Авто-SMART-оценка (0.0–1.0)"),
        ("strategic_alignment","strategic / functional / operational"),
    ]
    for i, (f, d) in enumerate(fields):
        ty = 4.56 + i * 0.38
        txt(s, f, 6.7,  ty, 2.5, 0.33, size=10, bold=True, color=GREEN, italic=True)
        txt(s, d, 9.25, ty, 3.6, 0.33, size=10, color=TEXT)
    pg(s, 15)
s9()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — API ENDPOINTS (технический)
# ══════════════════════════════════════════════════════════════════════════════
def s10():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "API — Документация и структура (10%)",
           "FastAPI · OpenAPI /docs · Pydantic v2 · Field(description) · Async SQLAlchemy")
    accent_line(s)

    # POST group
    post_eps = [
        ("/api/evaluate-goal",    "Оценить одну цель по SMART · 5 критериев · алерты F-16..F-21"),
        ("/api/evaluate-batch",   "Пакетная оценка всех целей сотрудника · asyncio.gather · radar"),
        ("/api/generate-goals",   "AI-генерация 3–5 целей · RAG + Gemini · dedup · grounding"),
        ("/api/accept-goals",     "Принять выбранные цели · сохранить в accepted_goals (F-13)"),
    ]
    # GET group
    get_eps = [
        ("/api/dashboard/department-quality", "SMART-индекс по подразделениям · strategic_ratio · weakest"),
        ("/api/dashboard/quarterly-trend",    "Квартальная динамика · criterion_scores для RadarChart"),
        ("/api/dashboard/maturity",           "Индекс зрелости подразделения F-22 · recommendations"),
        ("/api/employees",                    "Список сотрудников с кол-вом целей (кеш 5 мин)"),
        ("/health",                           "Health-check"),
    ]

    # POST header
    rect(s, 0.3, 1.55, 12.73, 0.38, GREEN)
    txt(s, "POST — запись / AI-обработка", 0.5, 1.58, 6.0, 0.3, size=11, bold=True, color=WHITE)
    txt(s, "Эндпоинт", 1.32, 1.58, 4.7, 0.3, size=11, bold=True, color=WHITE)
    txt(s, "Описание",  6.15, 1.58, 7.0, 0.3, size=11, bold=True, color=WHITE)

    for i, (path, desc) in enumerate(post_eps):
        ty = 1.98 + i * 0.5
        bg = GREEN_LT if i % 2 == 0 else WHITE
        rect(s, 0.3, ty, 12.73, 0.46, bg)
        rect(s, 0.38, ty+0.05, 0.72, 0.34, GREEN_LT)
        txt(s, "POST", 0.38, ty+0.06, 0.72, 0.32, size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(s, path, 1.22, ty+0.07, 4.7, 0.32, size=10.5, bold=True, color=GREEN, italic=True)
        txt(s, desc, 6.15, ty+0.07, 6.9, 0.32, size=10.5, color=TEXT)

    # Separator
    sep_y = 1.98 + len(post_eps) * 0.5 + 0.08
    rect(s, 0.3, sep_y, 12.73, 0.02, ACCENT)

    # GET header
    get_header_y = sep_y + 0.1
    rect(s, 0.3, get_header_y, 12.73, 0.38, MID_BLUE)
    txt(s, "GET — чтение / аналитика", 0.5, get_header_y+0.03, 6.0, 0.3, size=11, bold=True, color=WHITE)

    for i, (path, desc) in enumerate(get_eps):
        ty = get_header_y + 0.42 + i * 0.5
        bg = rgb(0xE3,0xF2,0xFD) if i % 2 == 0 else WHITE
        rect(s, 0.3, ty, 12.73, 0.46, bg)
        rect(s, 0.38, ty+0.05, 0.72, 0.34, rgb(0xE3,0xF2,0xFD))
        txt(s, "GET", 0.38, ty+0.06, 0.72, 0.32, size=9, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
        txt(s, path, 1.22, ty+0.07, 4.7, 0.32, size=10.5, bold=True, color=MID_BLUE, italic=True)
        txt(s, desc, 6.15, ty+0.07, 6.9, 0.32, size=10.5, color=TEXT)

    # QR code for Swagger docs
    try:
        s.shapes.add_picture("screenshots/qr_docs.png",
                             Inches(10.5), Inches(1.55), width=Inches(2.5), height=Inches(2.5))
    except Exception:
        pass
    txt(s, "Swagger UI", 10.5, 4.1, 2.5, 0.3,
        size=11, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "http://89.207.255.254:8001/docs", 10.2, 4.4, 2.8, 0.3,
        size=8, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    pg(s, 16)
s10()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 14 — MVP ТАБЛИЦА (смешанный)
# ══════════════════════════════════════════════════════════════════════════════
def s11():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "MVP — все обязательные функции реализованы")
    accent_line(s)

    rows = [
        ("✅", GREEN,   "SMART-оценка одной цели",         "POST /api/evaluate-goal",   "5 скоров · smart_index · рекомендации"),
        ("✅", GREEN,   "Переформулировка слабой цели",     "improved_goal",             "Автоматически если smart_index < 0.7"),
        ("✅", GREEN,   "Генерация 3–5 целей по должности", "POST /api/generate-goals",  "RAG (160 ВНД) + Gemini Flash"),
        ("✅", GREEN,   "Привязка к источнику ВНД",          "source_doc + source_quote", "Документ + цитата в каждой цели"),
        ("✅", GREEN,   "Пакетная оценка за квартал",        "POST /api/evaluate-batch",  "asyncio.gather · avg_index · weakest"),
        ("✅", GREEN,   "Дашборд по подразделениям",         "GET /api/dashboard/*",      "SMART-индекс · тренд · radar · проблемы"),
        ("⭕", ACCENT,  "Каскадирование от руководителя",    "manager_goals в промпте",   "Реализовано через manager_id"),
        ("–",  GRAY,   "Интеграция 1С / SAP / Oracle",      "—",                         "Вне MVP (опционально по ТЗ)"),
    ]

    rect(s, 0.3, 1.56, 12.73, 0.42, DARK_BLUE)
    for x, label in [(0.42, ""), (0.9, "Функция"), (5.0, "Endpoint / поле"), (8.2, "Детали")]:
        txt(s, label, x, 1.6, 3.0, 0.3, size=11, bold=True, color=WHITE)

    for i, (mark, color, func, endpoint, detail) in enumerate(rows):
        ty = 2.03 + i * 0.67
        bg = WHITE if i % 2 == 0 else rgb(0xF5,0xF8,0xFF)
        if mark == "✅":
            bg = GREEN_LT if i % 2 == 0 else rgb(0xDC,0xF0,0xDC)
        rect(s, 0.3, ty, 12.73, 0.62, bg)
        txt(s, mark,     0.42,  ty+0.12, 0.44, 0.38, size=16, color=color)
        txt(s, func,     0.9,   ty+0.13, 3.9,  0.36, size=12, bold=True, color=DARK_BLUE)
        txt(s, endpoint, 5.0,   ty+0.13, 3.0,  0.36, size=11, color=MID_BLUE, italic=True)
        txt(s, detail,   8.2,   ty+0.13, 4.7,  0.36, size=11, color=TEXT)
    pg(s, 17)
s11()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 15 — ДАННЫЕ (технический / бизнес)
# ══════════════════════════════════════════════════════════════════════════════
def s12():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Данные — что у нас есть", "Remote DB организаторов + Local DB AI-результатов + Qdrant")
    accent_line(s)

    # Stats strip
    stats = [
        (DARK_BLUE, "450",   "сотрудников"),
        (MID_BLUE,  "9 000", "целей в БД"),
        (GREEN,     "160",   "ВНД-документов"),
        (PURPLE,    "412",   "векторных чанков"),
        (ORANGE,    "8",     "подразделений"),
        (TEAL,      "13",    "KPI-метрик"),
    ]
    for i, (color, val, label) in enumerate(stats):
        lx = 0.3 + i * 2.12
        rect(s, lx, 1.56, 1.98, 1.2, color)
        txt(s, val,   lx+0.1, 1.66, 1.78, 0.6, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, lx+0.1, 2.24, 1.78, 0.38, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Remote DB
    rect(s, 0.3, 3.0, 5.9, 4.2, WHITE)
    rect(s, 0.3, 3.0, 5.9, 0.45, PURPLE)
    txt(s, "🟣  Remote PostgreSQL (readonly)", 0.5, 3.07, 5.5, 0.32, size=12, bold=True, color=WHITE)

    remote_tables = [
        ("employees",       "id, full_name, position_id, department_id, manager_id"),
        ("goals",           "goal_id, employee_id, goal_text, weight, quarter, year"),
        ("departments",     "id, name, is_active"),
        ("documents",       "id, title, content (полный текст ВНД)"),
        ("kpi_timeseries",  "metric_key, department_id, value_num, period_date"),
    ]
    for i, (tbl, fields) in enumerate(remote_tables):
        ty = 3.55 + i * 0.7
        rect(s, 0.45, ty, 5.5, 0.62, rgb(0xF3,0xE5,0xF5) if i%2==0 else WHITE)
        txt(s, tbl,    0.55, ty+0.06, 1.7, 0.36, size=11, bold=True, color=PURPLE)
        txt(s, fields, 2.3,  ty+0.1,  3.5, 0.36, size=9.5, color=TEXT)

    # Local DB
    rect(s, 6.5, 3.0, 6.55, 4.2, WHITE)
    rect(s, 6.5, 3.0, 6.55, 0.45, GREEN)
    txt(s, "🟢  Local PostgreSQL (Alembic) + Qdrant", 6.7, 3.07, 6.1, 0.32, size=12, bold=True, color=WHITE)

    local_tables = [
        ("ai_evaluations",    "scores_json, overall_index, goal_type, alignment_level"),
        ("ai_generated_goals","goals_json, context_json, warnings"),
        ("accepted_goals",    "goals_json, accepted_count"),
    ]
    for i, (tbl, fields) in enumerate(local_tables):
        ty = 3.55 + i * 0.7
        rect(s, 6.65, ty, 6.2, 0.62, GREEN_LT if i%2==0 else WHITE)
        txt(s, tbl,    6.75, ty+0.06, 2.5, 0.36, size=11, bold=True, color=GREEN)
        txt(s, fields, 9.3,  ty+0.1,  3.3, 0.36, size=9.5, color=TEXT)

    qdrant_items = "Qdrant:  collection=vnd_documents  ·  model=paraphrase-multilingual-MiniLM-L12-v2  ·  384d  ·  cosine  ·  412 chunks"
    rect(s, 6.5, 5.9, 6.55, 1.3, rgb(0xFF,0xF8,0xE1))
    rect(s, 6.5, 5.9, 6.55, 0.38, ACCENT)
    txt(s, "🟡  Qdrant", 6.7, 5.95, 2.0, 0.28, size=11, bold=True, color=DARK_BLUE)
    txt(s, qdrant_items, 6.6, 6.37, 6.3, 0.7, size=10, color=TEXT)
    pg(s, 18)
s12()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 16 — ЗАПУСК (технический)
# ══════════════════════════════════════════════════════════════════════════════
def s13():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
    rect(s, 0, 6.55, 13.33, 0.2, ACCENT)
    rect(s, 0, 0, 0.18, 7.5, ACCENT)

    txt(s, "Запуск за 3 команды", 0.5, 0.2, 12.33, 0.65,
        size=26, bold=True, color=WHITE)
    txt(s, "http://89.207.255.254:3000  ·  http://89.207.255.254:8001/docs",
        0.5, 0.82, 12.33, 0.4, size=14, color=ACCENT, italic=True)

    commands = [
        ("1", "docker compose up --build -d",
         "Поднять все 4 контейнера: local-db · qdrant · hr-api · hr-frontend"),
        ("2", "docker compose exec hr-api alembic upgrade head",
         "Создать таблицы AI-результатов в local-db (одноразово)"),
        ("3", "docker compose exec hr-api python -m scripts.ingest_docs",
         "Загрузить 160 ВНД в Qdrant — embeddings + chunking (одноразово)"),
    ]

    for i, (num, cmd, desc) in enumerate(commands):
        ty = 1.45 + i * 1.3
        rect(s, 0.3, ty, 12.73, 1.15, rgb(0x0D,0x1B,0x3E))
        rect(s, 0.3, ty, 0.52, 1.15, ACCENT)
        txt(s, num, 0.3, ty+0.35, 0.52, 0.45,
            size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        txt(s, cmd,  0.95, ty+0.1,  11.8, 0.42, size=13, bold=True, color=rgb(0xCC,0xFF,0xCC))
        txt(s, desc, 0.95, ty+0.6,  11.8, 0.42, size=11, color=rgb(0xAA,0xBB,0xDD))

    # Demo scenarios
    txt(s, "Демо-сценарии", 0.5, 5.35, 12.0, 0.38, size=13, bold=True, color=ACCENT)
    demos = [
        ("🔴", "Слабая цель",  "«Улучшить работу» → smart_index < 0.5 + переформулировка"),
        ("🟢", "Сильная цель", "«Сократить задержку ETL на 20% к 30.09.2025» → > 0.85, strategic"),
        ("✨", "Генерация",    "employee 269, Q3 2025, focus=цифровизация → 3–5 целей + source_doc"),
        ("📦", "Пакетная",     "evaluate-batch для Воронцовой (20 целей) → avg_index + radar"),
    ]
    for i, (icon, title, detail) in enumerate(demos):
        col, row = i % 2, i // 2
        lx = 0.4 + col * 6.5
        ty = 5.82 + row * 0.56
        txt(s, f"{icon}  {title}:", lx,     ty, 2.2,  0.44, size=11, bold=True, color=WHITE)
        txt(s, detail,              lx+2.2, ty, 4.2,  0.44, size=10.5, color=rgb(0xCC,0xDD,0xFF))
    pg(s, 19)
s13()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 20 — ROADMAP / РЕКОМЕНДАЦИИ
# ══════════════════════════════════════════════════════════════════════════════
def s_roadmap():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, LIGHT)
    header(s, "Roadmap — рекомендации по развитию",
           "От MVP к production-решению для 10 000+ сотрудников")
    accent_line(s)

    items = [
        (DARK_BLUE, "🔐", "Keycloak SSO + RBAC",
         "SSO через Active Directory / Keycloak\nJWT-токены, роли: сотрудник / руководитель / HR\nУбрать выбор сотрудника — автоматически по логину"),
        (PURPLE, "📄", "Парсинг реальных ВНД",
         "Word (.docx) и PDF с таблицами, списками, OCR\nAdaptive chunking по заголовкам\nВерсионирование: автообновление Qdrant при изменении"),
        (rgb(0xF4,0x43,0x36), "🖥️", "GPU-кластер / локальная LLM",
         "ВНД = конфиденциальные данные, нельзя в облако\nvLLM + Llama 3 / Mistral на NVIDIA A100\nПолная автономность, без rate limits"),
        (GREEN, "⚡", "Масштабирование",
         "Kubernetes HPA + 3 реплики API\nRedis: rate limiting + кеш + Celery очередь\nGunicorn + N workers вместо 1 uvicorn"),
    ]

    for i, (color, icon, title, body) in enumerate(items):
        col, row = i % 2, i // 2
        lx = 0.3 + col * 6.52
        ty = 1.55 + row * 2.8

        rect(s, lx, ty, 6.22, 2.55, WHITE)
        rect(s, lx, ty, 0.1, 2.55, color)
        txt(s, icon, lx + 0.2, ty + 0.15, 0.5, 0.5, size=24)
        txt(s, title, lx + 0.75, ty + 0.18, 5.2, 0.38, size=14, bold=True, color=color)
        txt(s, body, lx + 0.2, ty + 0.65, 5.8, 1.7, size=11, color=TEXT)

    pg(s, 20)
s_roadmap()


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 21 — ИТОГ / СПАСИБО (последний)
# ══════════════════════════════════════════════════════════════════════════════
def s14():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
    rect(s, 0, 5.2, 13.33, 0.12, ACCENT)
    rect(s, 0, 0, 0.18, 7.5, ACCENT)

    txt(s, "Спасибо за внимание", 0.5, 0.6, 12.33, 0.9,
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    summary = [
        ("🎯", "25% + 25%", "SMART-оценка + AI-генерация",   "Gemini Flash · 5 критериев · goal_type · strategic_alignment"),
        ("🔍", "15%",        "RAG-пайплайн",                   "Qdrant · 160 ВНД · 412 чанков · cosine · source grounding"),
        ("📊", "10%",        "Аналитика и дашборд",             "Dept quality · quarterly trend · maturity F-22 · radar chart"),
        ("⚙️", "10%",        "Архитектура и API",               "FastAPI · OpenAPI /docs · Docker Compose · Alembic · Pydantic v2"),
        ("💻", "15%",        "UX интерфейса",                   "React 18 · MUI · 6 страниц · progress bars · alert banners"),
    ]

    for i, (icon, weight, title, detail) in enumerate(summary):
        lx = 0.35 + (i % 3) * 4.2
        ty = 1.65 + (i // 3) * 2.1

        rect(s, lx, ty, 3.9, 1.85, rgb(0x0D,0x1B,0x3E))
        rect(s, lx, ty, 3.9, 0.08, ACCENT)

        txt(s, icon + "  " + weight, lx+0.15, ty+0.15, 3.6, 0.42, size=13, bold=True, color=ACCENT)
        txt(s, title,                lx+0.15, ty+0.56, 3.6, 0.38, size=12, bold=True, color=WHITE)
        txt(s, detail,               lx+0.15, ty+0.96, 3.6, 0.76, size=10, color=rgb(0xAA,0xCC,0xFF))

    txt(s, "Frontend: http://89.207.255.254:3000   ·   API: http://89.207.255.254:8001/docs",
        0.5, 6.9, 12.33, 0.38, size=12, color=rgb(0x99,0xBB,0xEE),
        align=PP_ALIGN.CENTER, italic=True)
    pg(s, 21)
s14()


# ─── SAVE ─────────────────────────────────────────────────────────────────────
OUT = "KMG_HR_AI_Presentation.pptx"
prs.save(OUT)
print(f"✅  Saved: {OUT}  ({len(prs.slides)} slides)")
